""".pptx — python-pptx 遍历 slide.shapes，按坐标排序模拟阅读顺序
.ppt  — LibreOffice 转 .pptx 后解析
"""

import os
import subprocess
import tempfile
from typing import List

import numpy as np
from PIL import Image
from io import BytesIO

from .base import ContentElement, DocumentLoader
from .ocr import ocr_image

SHAPE_PICTURE = 13
SHAPE_GROUP = 6


class PptxLoader(DocumentLoader):
    """.pptx 加载器：文本框/表格/图片/组合形状 → ContentElement。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def load(self, file_path: str) -> List[ContentElement]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        elements: List[ContentElement] = []
        file_name = os.path.basename(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            # 按坐标排序：从上到下、从左到右
            sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
            for shape in sorted_shapes:
                self._extract_shape(shape, file_path, file_name, slide_num, elements)

        return elements

    def _extract_shape(self, shape, file_path: str, file_name: str,
                       slide_num: int, elements: List[ContentElement]):
        # 文本框
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    elements.append(ContentElement(
                        content=text,
                        doc_type="paragraph",
                        metadata={"file_path": file_path, "file_name": file_name,
                                  "slide": slide_num},
                    ))

        # 表格
        if shape.has_table:
            for row in shape.table.rows:
                cells = [
                    " ".join(p.text.strip() for p in cell.text_frame.paragraphs
                             if p.text.strip())
                    for cell in row.cells
                ]
                row_text = " | ".join(cells)
                if row_text.strip():
                    elements.append(ContentElement(
                        content=row_text,
                        doc_type="table_row",
                        metadata={"file_path": file_path, "file_name": file_name,
                                  "slide": slide_num},
                    ))

        # 图片
        if shape.shape_type == SHAPE_PICTURE:
            try:
                pil_img = Image.open(BytesIO(shape.image.blob))
                ocr_text = ocr_image(np.array(pil_img))
                if ocr_text:
                    elements.append(ContentElement(
                        content=ocr_text,
                        doc_type="image",
                        metadata={"file_path": file_path, "file_name": file_name,
                                  "slide": slide_num},
                    ))
            except Exception:
                pass

        # 组合形状 — 递归
        elif shape.shape_type == SHAPE_GROUP:
            for child in shape.shapes:
                self._extract_shape(child, file_path, file_name, slide_num, elements)


class PptLoader(DocumentLoader):
    """.ppt 旧格式加载器：LibreOffice 转 .pptx 后复用 PptxLoader。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".ppt"]

    def load(self, file_path: str) -> List[ContentElement]:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pptx",
                     "--outdir", tmpdir, file_path],
                    capture_output=True, timeout=60,
                )
                pptx_files = [f for f in os.listdir(tmpdir) if f.endswith(".pptx")]
                if pptx_files:
                    return PptxLoader().load(os.path.join(tmpdir, pptx_files[0]))
        except FileNotFoundError:
            raise RuntimeError(
                f"无法解析 .ppt 文件: {file_path}。"
                "请安装 LibreOffice: brew install libreoffice"
            )
        return []
